#!/usr/bin/env python3
"""
PowerPoint Certificate Extractor

Extracts text positions, sizes, and formatting from PowerPoint certificate templates.
This provides exact positioning data from the original design source.
"""
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
import json

def extract_pptx_info(pptx_path):
    """Extract text box positions and properties from PPTX."""
    prs = Presentation(pptx_path)
    
    info = {
        'file': str(pptx_path),
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            'slide_number': slide_idx + 1,
            'width': prs.slide_width.pt,
            'height': prs.slide_height.pt,
            'text_boxes': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Get position and size
                left = shape.left.pt if hasattr(shape.left, 'pt') else 0
                top = shape.top.pt if hasattr(shape.top, 'pt') else 0
                width = shape.width.pt if hasattr(shape.width, 'pt') else 0
                height = shape.height.pt if hasattr(shape.height, 'pt') else 0
                
                # Calculate center position (normalized)
                center_x = (left + width / 2) / slide_info['width']
                center_y = (top + height / 2) / slide_info['height']
                
                # Get text properties
                text_box = {
                    'text': shape.text,
                    'left_pt': left,
                    'top_pt': top,
                    'width_pt': width,
                    'height_pt': height,
                    'center_x': center_x,
                    'center_y': center_y,
                    'center_x_pt': left + width / 2,
                    'center_y_pt': top + height / 2
                }
                
                # Try to get font information
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                text_box['font_size_pt'] = run.font.size.pt
                            if run.font.name:
                                text_box['font_name'] = run.font.name
                            if run.font.bold is not None:
                                text_box['bold'] = run.font.bold
                            break
                        break
                
                slide_info['text_boxes'].append(text_box)
        
        info['slides'].append(slide_info)
    
    return info


def match_fields_to_template(sample_info, template_info):
    """Match text fields from sample to template structure."""
    # Assuming first slide
    sample_boxes = sample_info['slides'][0]['text_boxes'] if sample_info['slides'] else []
    
    # Map common patterns
    field_mapping = {}
    
    for box in sample_boxes:
        text = box['text'].upper().strip()
        
        # Identify fields by content
        if 'SAMPLE' in text and 'NAME' in text:
            field_mapping['name'] = box
        elif 'SAMPLE' in text and 'EVENT' in text:
            field_mapping['event'] = box
        elif 'SAMPLE' in text and 'ORG' in text:
            field_mapping['organiser'] = box
        # Generic patterns
        elif 'NAME' in text:
            field_mapping.setdefault('name', box)
        elif 'EVENT' in text:
            field_mapping.setdefault('event', box)
        elif 'ORGANIS' in text or 'ORGANIZ' in text:
            field_mapping.setdefault('organiser', box)
    
    return field_mapping


def create_offsets_from_pptx(sample_pptx, template_pptx=None, output_path=None):
    """Create offset configuration from PPTX files."""
    print("=" * 80)
    print("EXTRACTING POSITIONS FROM POWERPOINT")
    print("=" * 80)
    print()
    
    # Extract info
    print(f"Reading: {sample_pptx}")
    sample_info = extract_pptx_info(sample_pptx)
    
    print(f"  Slide size: {sample_info['slides'][0]['width']}pt x {sample_info['slides'][0]['height']}pt")
    print(f"  Text boxes found: {len(sample_info['slides'][0]['text_boxes'])}")
    print()
    
    # Match fields
    field_mapping = match_fields_to_template(sample_info, None)
    
    print("Detected fields:")
    for field_name, box_info in field_mapping.items():
        print(f"\n  {field_name.upper()}:")
        print(f"    Text: '{box_info['text']}'")
        print(f"    Position: ({box_info['center_x']:.6f}, {box_info['center_y']:.6f})")
        print(f"    Position (pt): ({box_info['center_x_pt']:.1f}, {box_info['center_y_pt']:.1f})")
        print(f"    Size: {box_info['width_pt']:.1f} x {box_info['height_pt']:.1f} pt")
        if 'font_size_pt' in box_info:
            print(f"    Font size: {box_info['font_size_pt']:.1f}pt")
        if 'font_name' in box_info:
            print(f"    Font: {box_info['font_name']}")
    
    # Create configuration
    config = {
        'version': '7.0',
        'source': 'powerpoint_extraction',
        'extraction_date': '2025-11-01',
        'description': 'Field positions extracted directly from PowerPoint PPTX file',
        'slide_dimensions': {
            'width_pt': sample_info['slides'][0]['width'],
            'height_pt': sample_info['slides'][0]['height']
        },
        'fields': {}
    }
    
    for field_name, box_info in field_mapping.items():
        config['fields'][field_name] = {
            'x': box_info['center_x'],
            'y': box_info['center_y'],
            'font_size': int(box_info.get('font_size_pt', 70)),
            'baseline_offset': 0,
            'original_text': box_info['text'],
            'width_pt': box_info['width_pt'],
            'height_pt': box_info['height_pt']
        }
    
    # Save if output path provided
    if output_path:
        with open(output_path, 'w') as f:
            json.dump(config, f, indent=2)
        print(f"\n✅ Configuration saved to: {output_path}")
    
    return config


def main():
    """Main function."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Extract positions from PowerPoint')
    parser.add_argument('sample_pptx', help='Path to sample certificate PPTX')
    parser.add_argument('--template', help='Path to template PPTX (optional)')
    parser.add_argument('--output', help='Output JSON path', 
                       default='templates/pptx_extracted_offsets.json')
    
    args = parser.parse_args()
    
    if not Path(args.sample_pptx).exists():
        print(f"Error: File not found: {args.sample_pptx}")
        return 1
    
    try:
        config = create_offsets_from_pptx(args.sample_pptx, args.template, args.output)
        
        print("\n" + "=" * 80)
        print("EXTRACTION COMPLETE")
        print("=" * 80)
        print("\nUse this configuration with the certificate renderer.")
        print("The positions are extracted directly from PowerPoint's layout,")
        print("providing exact coordinates from the original design.")
        
        return 0
        
    except ImportError as e:
        print(f"\n❌ Error: {e}")
        print("\nPlease install python-pptx:")
        print("  pip install python-pptx")
        return 1
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
